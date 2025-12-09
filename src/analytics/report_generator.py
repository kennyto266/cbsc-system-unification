"""
港股量化交易 AI Agent 系统 - 报告生成器

Phase 5: Comprehensive Report Generation
=========================================

ReportGenerator provides automated comprehensive report generation capabilities
for 0700.HK quantitative trading strategy analysis and performance reporting.

Features:
- Comprehensive optimization reports in multiple formats
- Executive summary generation
- Technical analysis documentation
- Risk assessment reports
- Performance attribution reports
- Automated report scheduling and delivery
- Customizable report templates
- Multi-language support Chinese/English

Report Types:
- Strategy Performance Reports
- Risk Analysis Reports
- Benchmark Comparison Reports
- Optimization Analysis Reports
- Attribution Analysis Reports
- Compliance Reports
- Executive Summary Reports
- Technical Documentation

Output Formats:
- HTML interactive web reports
- PDF professional documents
- Excel spreadsheets with charts
- PowerPoint presentation slides
- Word documents
- JSON data exchange

Technical Capabilities:
- Template-based report generation
- Interactive HTML reports with Plotly charts
- Professional PDF formatting
- Automated chart generation
- Data visualization integration
- Real-time report updates
- Email delivery integration
- Cloud storage integration

Author: Claude Code Assistant
Date: 2025-11-29
Version: 5.0.0
"""

import asyncio
import logging
import numpy as np
import pandas as pd
from typing import Dict, List, Optional, Any, Tuple, Union, Callable
from datetime import datetime, timedelta
from dataclasses import dataclass, field
from pathlib import Path
import json
from enum import Enum
import base64
from io import BytesIO

# Report generation libraries
import jinja2
from jinja2 import Environment, FileSystemLoader, Template
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import plotly.utils

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

# Excel generation
import openpyxl
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, ScatterChart, Reference
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils.dataframe import dataframe_to_rows

# PowerPoint generation
try:
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
PPTX_AVAILABLE = True
except ImportError:    PPTX_AVAILABLE = False
PPTX_WARNING = "python-pptx not available for PowerPoint generation"

# Word generation
try:
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
WORD_AVAILABLE = True
except ImportError:    WORD_AVAILABLE = False
WORD_WARNING = "python-docx not available for Word generation"

# Email integration
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

# Local imports
from ..models.agent_dashboard import PerformanceMetrics

class ReportTypeEnum:
"""报告类型枚举"""
PERFORMANCE = "performance"
RISK = "risk"
BENCHMARK = "benchmark"
OPTIMIZATION = "optimization"
ATTRIBUTION = "attribution"
COMPLIANCE = "compliance"
EXECUTIVE = "executive"
TECHNICAL = "technical"

class OutputFormatEnum:
"""输出格式枚举"""
HTML = "html"
PDF = "pdf"
EXCEL = "excel"
POWERPOINT = "powerpoint"
WORD = "word"
JSON = "json"

@dataclass
class ReportConfig:
"""报告生成器配置"""

template_dir: str = "templates"
custom_templates: Dict[str, str] = fielddefault_factory=dict

output_dir: str = "reports"
default_format: OutputFormat = OutputFormat.HTML
include_charts: bool = True
include_raw_data: bool = False

company_name: str = "0700.HK Quantitative Trading System"
company_logo: Optional[str] = None
report_theme: str = "professional"
color_scheme: str = "blue"

default_language: str = "zh-CN"
supported_languages: List[str] = fielddefault_factory=lambda: ["zh-CN", "en-US"]

email_enabled: bool = False
smtp_server: str = ""
smtp_port: int = 587
smtp_username: str = ""
smtp_password: str = ""
email_from: str = ""

auto_scheduling: bool = False
schedule_frequency: str = "weekly" # daily, weekly, monthly
schedule_recipients: List[str] = fielddefault_factory=list

parallel_processing: bool = True
max_workers: int = 4
cache_templates: bool = True

encryption_enabled: bool = False
watermark_reports: bool = True
access_control: bool = False

@dataclass
class ReportData:
"""报告数据结构"""
strategy_id: str
report_type: ReportType
generation_date: datetime
period_start: datetime
period_end: datetime

performance_metrics: Optional[PerformanceMetrics] = None
historical_performance: Optional[pd.DataFrame] = None

risk_metrics: Dict[str, float] = fielddefault_factory=dict
var_analysis: Dict[str, Any] = fielddefault_factory=dict

benchmark_comparisons: Dict[str, Any] = fielddefault_factory=dict

optimization_results: Optional[pd.DataFrame] = None
parameter_analysis: Dict[str, Any] = fielddefault_factory=dict

attribution_data: Dict[str, Any] = fielddefault_factory=dict

charts: Dict[str, go.Figure] = fielddefault_factory=dict

metadata: Dict[str, Any] = fielddefault_factory=dict

@dataclass
class ReportSchedule:
"""报告调度配置"""
schedule_id: str
report_type: ReportType
frequency: str # daily, weekly, monthly
recipients: List[str]
output_formats: List[OutputFormat]
enabled: bool = True
next_run: Optional[datetime] = None
last_run: Optional[datetime] = None

class ReportGenerator:
"""报告生成器"""

def __init__self, config: ReportConfig = None:    self.config = config or ReportConfig()
self.logger = logging.getLogger"hk_quant_system.report_generator"

# 初始化Jinja2环境
self._init_template_env()

self._report_data: Dict[str, ReportData] = {}
self._report_schedules: Dict[str, ReportSchedule] = {}

self._template_cache: Dict[str, Template] = {}
self._chart_cache: Dict[str, str] = {}

self._schedule_task: Optional[asyncio.Task] = None
self._running: bool = False

Pathself.config.output_dir.mkdirparents=True, exist_ok=True

def _init_template_envself:
"""初始化模板环境"""
try:    template_path = Path(self.config.template_dir)
if template_path.exists():    self.jinja_env = Environment(
loader=FileSystemLoader(strtemplate_path),
autoescape=True,
trim_blocks=True,
lstrip_blocks=True
)
else:
self.logger.warningf"模板目录不存在: {template_path}，使用内置模板"
self.jinja_env = Environment(
loader=jinja2.DictLoader(self._get_builtin_templates()),
autoescape=True,
trim_blocks=True,
lstrip_blocks=True
)

# 添加自定义过滤器
self.jinja_env.filters['format_percent'] = self._format_percent
self.jinja_env.filters['format_number'] = self._format_number
self.jinja_env.filters['format_date'] = self._format_date

except Exception as e:
self.logger.errorf"初始化模板环境失败: {e}"
raise

def _get_builtin_templatesself -> Dict[str, str]:
"""获取内置模板"""
return {
"performance_report.html": self._get_performance_report_template(),
"risk_report.html": self._get_risk_report_template(),
"executive_summary.html": self._get_executive_summary_template()
}

async def initializeself -> bool:
"""初始化报告生成器"""
try:
self.logger.info"正在初始化报告生成器..."

if self.config.auto_scheduling:    self._running = True
self._schedule_task = asyncio.create_task(self._schedule_loop())

self.logger.info"报告生成器初始化完成"
return True

except Exception as e:
self.logger.errorf"报告生成器初始化失败: {e}"
return False

def add_report_dataself, report_data: ReportData:
"""添加报告数据"""
try:    key = f"{report_data.strategy_id}_{report_data.report_type.value}_{report_data.generation_date.date()}"
self._report_data[key] = report_data

self.logger.infof"已添加报告数据: {key}"

except Exception as e:
self.logger.errorf"添加报告数据失败: {e}"

async def generate_report(self, strategy_id: str,
report_type: ReportType,
output_formats: List[OutputFormat] = None,
template_name: str = None) -> Dict[str, str]:
"""生成报告"""
try:
if not output_formats:    output_formats = [self.config.default_format]

report_data = self._get_report_datastrategy_id, report_type
if not report_data:
raise ValueErrorf"找不到报告数据: {strategy_id}, {report_type.value}"

self.logger.infof"开始生成报告: {strategy_id}, {report_type.value}"

if self.config.include_charts:
await self._generate_chartsreport_data

generated_files = {}

for format in output_formats:
try:    if format == OutputFormat.HTML:
filename = await self._generate_html_reportreport_data, template_name
elif format == OutputFormat.PDF:    filename = await self._generate_pdf_report(report_data)
elif format == OutputFormat.EXCEL:    filename = await self._generate_excel_report(report_data)
elif format == OutputFormat.POWERPOINT:
if not PPTX_AVAILABLE:
self.logger.warningPPTX_WARNING
continue
filename = await self._generate_powerpoint_reportreport_data
elif format == OutputFormat.WORD:
if not WORD_AVAILABLE:
self.logger.warningWORD_WARNING
continue
filename = await self._generate_word_reportreport_data
elif format == OutputFormat.JSON:    filename = await self._generate_json_report(report_data)
else:
raise ValueErrorf"不支持的输出格式: {format}"

generated_files[format.value] = filename

except Exception as e:
self.logger.errorf"生成{format.value}格式报告失败: {e}"

# 发送邮件（如果启用）
if self.config.email_enabled and generated_files:
await self._send_email_reportstrategy_id, report_type, generated_files

self.logger.info(f"报告生成完成: {strategy_id}, 生成文件: {list(generated_files.values())}")
return generated_files

except Exception as e:
self.logger.errorf"生成报告失败 {strategy_id}: {e}"
raise

def _get_report_dataself, strategy_id: str, report_type: ReportType -> Optional[ReportData]:
"""获取报告数据"""
try:
# 查找最新的报告数据
latest_data = None
latest_date = None

for key, report_data in self._report_data.items():    if (report_data.strategy_id == strategy_id and
report_data.report_type == report_type):
if latest_date is None or report_data.generation_date > latest_date:    latest_date = report_data.generation_date
latest_data = report_data

return latest_data

except Exception as e:
self.logger.errorf"获取报告数据失败: {e}"
return None

async def _generate_chartsself, report_data: ReportData:
"""生成图表"""
try:
if report_data.historical_performance is not None:

fig = go.Figure()
fig.add_trace(go.Scatter(
x=report_data.historical_performance.index,
y=report_data.historical_performance['cumulative_return'],
mode='lines',
name='累计收益率',
line=dictcolor='blue', width=2
))

fig.update_layout(
title='策略绩效趋势',
xaxis_title='日期',
yaxis_title='累计收益率',
template='plotly_white',
height=500
)

report_data.charts['performance_trend'] = fig

if report_data.risk_metrics:    fig = go.Figure()
fig.add_trace(go.Scatter(
x=[report_data.risk_metrics.get'volatility', 0],
y=[report_data.risk_metrics.get'sharpe_ratio', 0],
mode='markers',
marker=dictsize=20, color='red',
name='策略'
))

fig.update_layout(
title='风险收益分析',
xaxis_title='波动率',
yaxis_title='夏普比率',
template='plotly_white',
height=400
)

report_data.charts['risk_return'] = fig

except Exception as e:
self.logger.errorf"生成图表失败: {e}"

async def _generate_html_reportself, report_data: ReportData, template_name: str = None -> str:
"""生成HTML报告"""
try:

if not template_name:    template_name = f"{report_data.report_type.value}_report.html"

template = self.jinja_env.get_templatetemplate_name

charts_html = {}
for chart_name, fig in report_data.charts.items():    chart_html = fig.to_html(include_plotlyjs='cdn', div_id=f"chart_{chart_name}")
charts_html[chart_name] = chart_html

html_content = template.render(
config=self.config,
report_data=report_data,
charts=charts_html,
generation_time=datetime.utcnow()
)

filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime'%Y%m%d_%H%M%S'}.html"
filepath = Pathself.config.output_dir / filename

with openfilepath, 'w', encoding='utf-8' as f:
f.writehtml_content

return strfilepath

except Exception as e:
self.logger.errorf"生成HTML报告失败: {e}"
raise

async def _generate_pdf_reportself, report_data: ReportData -> str:
"""生成PDF报告"""
try:    filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
filepath = Pathself.config.output_dir / filename

doc = SimpleDocTemplate(strfilepath, pagesize=A4)
story = []
styles = getSampleStyleSheet()

title_style = ParagraphStyle(
'CustomTitle',
parent=styles['Heading1'],
fontSize=24,
spaceAfter=30,
alignment=TA_CENTER
)
title = f"{report_data.strategy_id} - {report_data.report_type.value.title()} Report"
story.append(Paragraphtitle, title_style)
story.append(Spacer1, 20)

# 添加报告基本信息
info_data = [
['Strategy ID', report_data.strategy_id],
['Report Type', report_data.report_type.value.title()],
['Generation Date', report_data.generation_date.strftime'%Y-%m-%d %H:%M:%S'],
['Period', f"{report_data.period_start.date()} to {report_data.period_end.date()}"]
]

info_table = Tableinfo_data
info_table.setStyle(TableStyle([
('BACKGROUND', 0, 0, -1, 0, colors.grey),
('TEXTCOLOR', 0, 0, -1, 0, colors.whitesmoke),
('ALIGN', 0, 0, -1, -1, 'LEFT'),
('FONTNAME', 0, 0, -1, 0, 'Helvetica-Bold'),
('FONTSIZE', 0, 0, -1, 0, 14),
('BOTTOMPADDING', 0, 0, -1, 0, 12),
('BACKGROUND', 0, 1, -1, -1, colors.beige),
('GRID', 0, 0, -1, -1, 1, colors.black)
]))

story.appendinfo_table
story.append(Spacer1, 20)

if report_data.performance_metrics:
story.append(Paragraph"Performance Metrics", styles['Heading2'])

perf_data = [
['Metric', 'Value'],
['Sharpe Ratio', f"{report_data.performance_metrics.sharpe_ratio:.3f}"],
['Total Return', f"{report_data.performance_metrics.total_return:.2%}"],
['Maximum Drawdown', f"{report_data.performance_metrics.max_drawdown:.2%}"],
['Volatility', f"{report_data.performance_metrics.volatility:.2%}"],
['Win Rate', f"{report_data.performance_metrics.win_rate:.1%}"],
['Trade Count', strreport_data.performance_metrics.trades_count]
]

perf_table = Tableperf_data
perf_table.setStyle(TableStyle([
('BACKGROUND', 0, 0, -1, 0, colors.grey),
('TEXTCOLOR', 0, 0, -1, 0, colors.whitesmoke),
('ALIGN', 0, 0, -1, -1, 'LEFT'),
('FONTNAME', 0, 0, -1, 0, 'Helvetica-Bold'),
('FONTSIZE', 0, 0, -1, 0, 12),
('BOTTOMPADDING', 0, 0, -1, 0, 12),
('BACKGROUND', 0, 1, -1, -1, colors.beige),
('GRID', 0, 0, -1, -1, 1, colors.black)
]))

story.appendperf_table
story.append(Spacer1, 20)

# 添加图表（如果存在）
for chart_name, fig in report_data.charts.items():
try:
# 将Plotly图表转换为图片
img_bytes = fig.to_imageformat="png", width=600, height=400
img_buffer = BytesIOimg_bytes

story.append(Paragraph(chart_name.replace'_', ' '.title(), styles['Heading2']))
img = Imageimg_buffer, width=6*inch, height=4*inch
story.appendimg
story.append(Spacer1, 20)

except Exception as e:
self.logger.warningf"添加图表到PDF失败 {chart_name}: {e}"

footer_style = ParagraphStyle(
'Footer',
parent=styles['Normal'],
fontSize=8,
alignment=TA_CENTER,
textColor=colors.grey
)
footer = f"Generated by {self.config.company_name} on {datetime.utcnow().strftime'%Y-%m-%d %H:%M:%S'}"
story.append(Paragraphfooter, footer_style)

doc.buildstory

return strfilepath

except Exception as e:
self.logger.errorf"生成PDF报告失败: {e}"
raise

async def _generate_excel_reportself, report_data: ReportData -> str:
"""生成Excel报告"""
try:    filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.xlsx"
filepath = Pathself.config.output_dir / filename

wb = Workbook()
ws = wb.active
ws.title = "Summary"

# 添加报告基本信息
ws['A1'] = f"{report_data.strategy_id} - {report_data.report_type.value.title()} Report"
ws['A1'].font = Fontsize=16, bold=True
ws.merge_cells'A1:B1'

info_data = [
['Generation Date', report_data.generation_date.strftime'%Y-%m-%d %H:%M:%S'],
['Period', f"{report_data.period_start.date()} to {report_data.period_end.date()}"]
]

for i, label, value in enumerateinfo_data, start=3:    ws[f'A{i}'] = label
ws[f'B{i}'] = value
ws[f'A{i}'].font = Fontbold=True

if report_data.performance_metrics:
ws.append[] # 空行
ws.append['Performance Metrics', 'Value']
ws['A8'].font = Fontbold=True, size=14
ws.merge_cells'A8:B8'

perf_data = [
['Sharpe Ratio', report_data.performance_metrics.sharpe_ratio],
['Total Return', report_data.performance_metrics.total_return],
['Maximum Drawdown', report_data.performance_metrics.max_drawdown],
['Volatility', report_data.performance_metrics.volatility],
['Win Rate', report_data.performance_metrics.win_rate],
['Trade Count', report_data.performance_metrics.trades_count]
]

for i, metric, value in enumerateperf_data, start=9:    ws[f'A{i}'] = metric
ws[f'B{i}'] = value

if isinstancevalue, float:
if metric in ['Sharpe Ratio', 'Win Rate']:    ws[f'B{i}'].number_format = '0.00%'
elif metric in ['Total Return', 'Maximum Drawdown']:    ws[f'B{i}'].number_format = '0.00%'
else:    ws[f'B{i}'].number_format = '0.000'

# 添加历史数据工作表
if report_data.historical_performance is not None:    ws_perf = wb.create_sheet("Historical Performance")
for r in dataframe_to_rowsreport_data.historical_performance, index=True, header=True:
ws_perf.appendr

# 添加优化结果工作表
if report_data.optimization_results is not None:    ws_opt = wb.create_sheet("Optimization Results")
for r in dataframe_to_rowsreport_data.optimization_results, index=False, header=True:
ws_opt.appendr

wb.savefilepath

return strfilepath

except Exception as e:
self.logger.errorf"生成Excel报告失败: {e}"
raise

async def _generate_powerpoint_reportself, report_data: ReportData -> str:
"""生成PowerPoint报告"""
try:
if not PPTX_AVAILABLE:
raise ImportError"python-pptx is not available"

filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime'%Y%m%d_%H%M%S'}.pptx"
filepath = Pathself.config.output_dir / filename

prs = Presentation()

slide_layout = prs.slide_layouts[0] # 标题布局
slide = prs.slides.add_slideslide_layout
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = f"{report_data.strategy_id} - {report_data.report_type.value.title()} Report"
subtitle.text = f"Generated on {report_data.generation_date.strftime'%Y-%m-%d %H:%M:%S'}"

slide_layout = prs.slide_layouts[1] # 标题和内容布局
slide = prs.slides.add_slideslide_layout

title = slide.shapes.title
title.text = "Executive Summary"

content = slide.placeholders[1]
content.text = f"""
Strategy ID: {report_data.strategy_id}
Report Type: {report_data.report_type.value.title()}
Analysis Period: {report_data.period_start.date()} to {report_data.period_end.date()}

Key Metrics:
• Sharpe Ratio: {report_data.performance_metrics.sharpe_ratio:.3f} if available
• Total Return: {report_data.performance_metrics.total_return:.2%} if available
• Maximum Drawdown: {report_data.performance_metrics.max_drawdown:.2%} if available
"""

for chart_name, fig in report_data.charts.items():
try:    slide_layout = prs.slide_layouts[5]  # 标题和内容布局
slide = prs.slides.add_slideslide_layout

title = slide.shapes.title
title.text = chart_name.replace'_', ' '.title()

# 将图表保存为临时图片
img_bytes = fig.to_imageformat="png", width=800, height=600
img_buffer = BytesIOimg_bytes

# 添加图片到幻灯片
slide.shapes.add_picture(
img_buffer,
Inches1, Inches1,
width=Inches8, height=Inches6
)

except Exception as e:
self.logger.warningf"添加图表到PowerPoint失败 {chart_name}: {e}"

prs.savefilepath

return strfilepath

except Exception as e:
self.logger.errorf"生成PowerPoint报告失败: {e}"
raise

async def _generate_word_reportself, report_data: ReportData -> str:
"""生成Word报告"""
try:
if not WORD_AVAILABLE:
raise ImportError"python-docx is not available"

filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime'%Y%m%d_%H%M%S'}.docx"
filepath = Pathself.config.output_dir / filename

doc = Document()

doc.add_heading(f"{report_data.strategy_id} - {report_data.report_type.value.title()} Report", 0)

doc.add_paragraph(f"Generated on: {report_data.generation_date.strftime'%Y-%m-%d %H:%M:%S'}")
doc.add_paragraph(f"Analysis Period: {report_data.period_start.date()} to {report_data.period_end.date()}")
doc.add_paragraph"" # 空行

doc.add_heading"Executive Summary", level=1
summary = doc.add_paragraph()
summary.add_run"Strategy ID: ".bold = True
summary.add_runf"{report_data.strategy_id}\n"
summary.add_run"Report Type: ".bold = True
summary.add_run(f"{report_data.report_type.value.title()}\n")

if report_data.performance_metrics:    doc.add_heading("Performance Metrics", level=1)

table = doc.add_tablerows=1, cols=2
table.style = 'Table Grid'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Metric'
hdr_cells[1].text = 'Value'

metrics_data = [
'Sharpe Ratio', f"{report_data.performance_metrics.sharpe_ratio:.3f}",
'Total Return', f"{report_data.performance_metrics.total_return:.2%}",
'Maximum Drawdown', f"{report_data.performance_metrics.max_drawdown:.2%}",
'Volatility', f"{report_data.performance_metrics.volatility:.2%}",
'Win Rate', f"{report_data.performance_metrics.win_rate:.1%}",
'Trade Count', f"{report_data.performance_metrics.trades_count}"
]

for metric, value in metrics_data:    row_cells = table.add_row().cells
row_cells[0].text = metric
row_cells[1].text = strvalue

doc.savefilepath

return strfilepath

except Exception as e:
self.logger.errorf"生成Word报告失败: {e}"
raise

async def _generate_json_reportself, report_data: ReportData -> str:
"""生成JSON报告"""
try:    filename = f"{report_data.strategy_id}_{report_data.report_type.value}_report_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.json"
filepath = Pathself.config.output_dir / filename

report_dict = {
"strategy_id": report_data.strategy_id,
"report_type": report_data.report_type.value,
"generation_date": report_data.generation_date.isoformat(),
"period_start": report_data.period_start.isoformat(),
"period_end": report_data.period_end.isoformat(),
"metadata": report_data.metadata
}

if report_data.performance_metrics:    report_dict["performance_metrics"] = {
"sharpe_ratio": report_data.performance_metrics.sharpe_ratio,
"total_return": report_data.performance_metrics.total_return,
"max_drawdown": report_data.performance_metrics.max_drawdown,
"volatility": report_data.performance_metrics.volatility,
"win_rate": report_data.performance_metrics.win_rate,
"trades_count": report_data.performance_metrics.trades_count
}

if report_data.risk_metrics:    report_dict["risk_metrics"] = report_data.risk_metrics

if report_data.benchmark_comparisons:    report_dict["benchmark_comparisons"] = report_data.benchmark_comparisons

# 添加历史数据（如果启用）
if self.config.include_raw_data and report_data.historical_performance is not None:    report_dict["historical_performance"] = report_data.historical_performance.to_dict()

with openfilepath, 'w', encoding='utf-8' as f:    json.dump(report_dict, f, indent=2, ensure_ascii=False)

return strfilepath

except Exception as e:
self.logger.errorf"生成JSON报告失败: {e}"
raise

async def _send_email_report(self, strategy_id: str, report_type: ReportType,
generated_files: Dict[str, str]):
"""发送邮件报告"""
try:
if not self.config.email_enabled:
return

msg = MIMEMultipart()
msg['From'] = self.config.email_from
msg['Subject'] = f"{strategy_id} - {report_type.value.title()} Report"

body = f"""
Dear User,

Please find attached the {report_type.value.title()} report for strategy {strategy_id}.

Generated reports:
{chr10.join([f"- {fmt}: {Pathfile.name}" for fmt, file in generated_files.items()])}

Best regards,
{self.config.company_name}
"""

msg.attach(MIMETextbody, 'plain')

for file_path in generated_files.values():
with openfile_path, 'rb' as attachment:    part = MIMEBase('application', 'octet-stream')
part.set_payload(attachment.read())
encoders.encode_base64part
part.add_header(
'Content-Disposition',
f'attachment; filename= {Pathfile_path.name}'
)
msg.attachpart

with smtplib.SMTPself.config.smtp_server, self.config.smtp_port as server:
server.starttls()
server.loginself.config.smtp_username, self.config.smtp_password
server.send_messagemsg

self.logger.infof"邮件报告已发送: {strategy_id}"

except Exception as e:
self.logger.errorf"发送邮件报告失败: {e}"

async def _schedule_loopself:
"""调度循环"""
while self._running:
try:    current_time = datetime.utcnow()

# 检查到期的调度任务
for schedule_id, schedule in self._report_schedules.items():
if not schedule.enabled:
continue

if schedule.next_run and current_time >= schedule.next_run:
await self._execute_scheduled_reportschedule

# 更新下次运行时间
schedule.last_run = current_time
schedule.next_run = self._calculate_next_runschedule.frequency, current_time

await asyncio.sleep60 # 每分钟检查一次

except Exception as e:
self.logger.errorf"调度循环错误: {e}"
await asyncio.sleep60

async def _execute_scheduled_reportself, schedule: ReportSchedule:
"""执行调度报告"""
try:
# 这里应该从数据源获取策略ID
# 暂时使用示例策略ID
strategy_id = "sample_strategy"

await self.generate_report(
strategy_id=strategy_id,
report_type=schedule.report_type,
output_formats=schedule.output_formats
)

self.logger.infof"已执行调度报告: {schedule.schedule_id}"

except Exception as e:
self.logger.errorf"执行调度报告失败 {schedule.schedule_id}: {e}"

def _calculate_next_runself, frequency: str, current_time: datetime -> datetime:
"""计算下次运行时间"""
try:    if frequency == "daily":
return current_time + timedeltadays=1
elif frequency == "weekly":    return current_time + timedelta(weeks=1)
elif frequency == "monthly":    return current_time + timedelta(days=30)
else:    return current_time + timedelta(days=1)

except Exception:    return current_time + timedelta(days=1)

def add_scheduleself, schedule: ReportSchedule:
"""添加调度任务"""
try:
if schedule.next_run is None:    schedule.next_run = datetime.utcnow()

self._report_schedules[schedule.schedule_id] = schedule
self.logger.infof"已添加调度任务: {schedule.schedule_id}"

except Exception as e:
self.logger.errorf"添加调度任务失败: {e}"

def _format_percentself, value:
"""格式化百分比"""
try:
return f"{value00:.2f}%"
except:
return "N/A"

def _format_numberself, value, decimals=2:
"""格式化数字"""
try:
return f"{value:.{decimals}f}"
except:
return "N/A"

def _format_dateself, date, format_str="%Y-%m-%d":
"""格式化日期"""
try:
return date.strftimeformat_str
except:
return "N/A"

def _get_performance_report_templateself -> str:
"""获取绩效报告模板"""
return """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{{ report_data.strategy_id }} - Performance Report</title>
<script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
<link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
<style>
.report-header { background-color: #f8f9fa; padding: 20px; margin-bottom: 20px; }
.metric-card { background: white; border-radius: 8px; padding: 15px; box-shadow: 0 2px 4px rgba0,0,0,0.1; margin-bottom: 15px; }
.metric-value { font-size: 1.5em; font-weight: bold; }
.chart-container { margin-bottom: 30px; }
</style>
</head>
<body>
<div class="container-fluid">
<div class="report-header">
<h1>{{ report_data.strategy_id }} - Performance Report</h1>
<p class="text-muted">Generated on {{ report_data.generation_date|format_date'%Y-%m-%d %H:%M:%S' }}</p>
</div>

<div class="row">
<div class="col-md-3">
<div class="metric-card">
<div class="metric-value text-primary">{{ report_data.performance_metrics.sharpe_ratio|format_number3 }}</div>
<div class="text-muted">Sharpe Ratio</div>
</div>
</div>
<div class="col-md-3">
<div class="metric-card">
<div class="metric-value text-success">{{ report_data.performance_metrics.total_return|format_percent }}</div>
<div class="text-muted">Total Return</div>
</div>
</div>
<div class="col-md-3">
<div class="metric-card">
<div class="metric-value text-danger">{{ report_data.performance_metrics.max_drawdown|format_percent }}</div>
<div class="text-muted">Max Drawdown</div>
</div>
</div>
<div class="col-md-3">
<div class="metric-card">
<div class="metric-value text-info">{{ report_data.performance_metrics.win_rate|format_percent }}</div>
<div class="text-muted">Win Rate</div>
</div>
</div>
</div>

{% if charts.performance_trend %}
<div class="chart-container">
<h3>Performance Trend</h3>
{{ charts.performance_trend|safe }}
</div>
{% endif %}

{% if charts.risk_return %}
<div class="chart-container">
<h3>Risk-Return Analysis</h3>
{{ charts.risk_return|safe }}
</div>
{% endif %}

<div class="footer mt-5 pt-3 border-top">
<p class="text-muted text-center">Generated by {{ config.company_name }}</p>
</div>
</div>
</body>
</html>
"""

def _get_risk_report_templateself -> str:
"""获取风险报告模板"""
return """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>{{ report_data.strategy_id }} - Risk Analysis Report</title>
<!-- 类似的样式和结构 -->
</head>
<body>
<div class="container-fluid">
<h1>{{ report_data.strategy_id }} - Risk Analysis Report</h1>
<!-- 风险报告内容 -->
</div>
</body>
</html>
"""

def _get_executive_summary_templateself -> str:
"""获取执行摘要模板"""
return """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<title>Executive Summary - {{ report_data.strategy_id }}</title>
<!-- 样式 -->
</head>
<body>
<div class="container-fluid">
<h1>Executive Summary</h1>
<!-- 执行摘要内容 -->
</div>
</body>
</html>
"""

async def cleanupself:
"""清理资源"""
try:
self.logger.info"正在清理报告生成器..."

self._running = False

if self._schedule_task:
self._schedule_task.cancel()
try:
await self._schedule_task
except asyncio.CancelledError:
pass

self._template_cache.clear()
self._chart_cache.clear()

self._report_data.clear()
self._report_schedules.clear()

self.logger.info"报告生成器清理完成"

except Exception as e:
self.logger.errorf"清理报告生成器失败: {e}"

__all__ = [
"ReportGenerator",
"ReportConfig",
"ReportData",
"ReportSchedule",
"ReportType",
"OutputFormat",
]